#!/usr/bin/env python3
"""
generate_vcenter_report.py
--------------------------
Reads an RVTools OR Nutanix Collector XLSX export (auto-detected) and produces
a PowerPoint report with slides grouped by datacenter:
  - Divider slide per datacenter
  - ESXi Hosts slide: one row per host with CPU model, memory, VM count, etc.
  - VM Resource Summary slide: storage allocated/used, memory, vCPUs, p:v ratio

Usage:
    python generate_vcenter_report.py <input.xlsx> [output.pptx]

Requires:
    pip install openpyxl python-pptx
"""

import sys
from pathlib import Path
from collections import defaultdict

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BG    = RGBColor(0x1E, 0x29, 0x3B)
ACCENT     = RGBColor(0x00, 0x8C, 0xD7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF8)
MID_GRAY   = RGBColor(0xBD, 0xC3, 0xC7)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
ROW_ALT    = RGBColor(0xEA, 0xF4, 0xFB)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def mib_to_gib(mib):
    if not mib: return 0.0
    return round(float(mib) / 1024, 1)

def mib_to_tib(mib):
    if not mib: return 0.0
    return round(float(mib) / (1024 * 1024), 2)

def mb_to_gib(mb):
    if not mb: return 0.0
    return round(float(mb) / 1024, 1)

def shorten_host(name):
    return name.split(".")[0] if name else name

def ratio_str(phys, virt):
    if not phys: return "N/A"
    return f"1 : {virt / phys:.1f}"

def rows_from_sheet(wb, sheet_name):
    ws = wb[sheet_name]
    rows = list(ws.iter_rows(values_only=True))
    headers = {h: i for i, h in enumerate(rows[0])}
    return rows[1:], headers

def get(row, headers, *keys):
    """Return the first key found in headers, or None."""
    for k in keys:
        if k in headers:
            v = row[headers[k]]
            return v
    return None


# ---------------------------------------------------------------------------
# Format detection
# ---------------------------------------------------------------------------
def detect_format(wb):
    sheets = wb.sheetnames
    if "vHosts" in sheets:
        return "nutanix"
    if "vHost" in sheets:
        return "rvtools"
    raise ValueError(f"Unrecognised XLSX format. Sheets found: {sheets}")


# ---------------------------------------------------------------------------
# RVTools loader
# ---------------------------------------------------------------------------
RVTOOLS_HOST_KEYS = {
    "name":             ["vHostName",        "Host"],
    "datacenter":       ["vHostDatacenter",  "Datacenter"],
    "cluster":          ["vHostCluster",     "Cluster"],
    "cpu_model":        ["vHostCpuModel",    "CPU Model"],
    "cpu_mhz":          ["vHostCpuMhz",      "Speed"],
    "num_sockets":      ["vHostNumCpu",      "# CPU"],
    "cores_per_socket": ["vHostCoresPerCPU", "Cores per CPU"],
    "total_cores":      ["vHostNumCpuCores", "# Cores"],
    "memory_mb":        ["vHostMemorySize",  "# Memory"],
    "vms_total":        ["vHostVMsTotal",    "# VMs total"],
    "vms_on":           ["vHostVMs",         "# VMs"],
    "vcpus":            ["vHostvCPUs",       "# vCPUs"],
    "vram_mib":         ["vHostvRAM",        "vRAM"],
}

RVTOOLS_INFO_KEYS = {
    "name":            ["vInfoVMName",               "VM"],
    "template":        ["vInfoTemplate",             "Template"],
    "cpus":            ["vInfoCPUs",                 "CPUs"],
    "memory_mib":      ["vInfoMemory",               "Memory"],
    "provisioned_mib": ["vInfoProvisioned",          "Provisioned MiB"],
    "inuse_mib":       ["vInfoInUse",                "In Use MiB"],
    "datacenter":      ["vInfoDataCenter",           "Datacenter"],
    "cluster":         ["vInfoCluster",              "Cluster"],
    "host":            ["vInfoHost",                 "Host"],
}

def _resolve(headers, keys):
    for k in keys:
        if k in headers:
            return headers[k]
    raise KeyError(f"Column not found — tried: {keys}")

def load_rvtools(wb):
    # Hosts
    host_rows, hi = rows_from_sheet(wb, "vHost")
    idx_h = {field: _resolve(hi, aliases) for field, aliases in RVTOOLS_HOST_KEYS.items()}

    hosts = []
    for row in host_rows:
        if not row[idx_h["name"]]:
            continue
        hosts.append({f: row[idx_h[f]] or (0 if f not in ("name","datacenter","cluster","cpu_model") else "")
                      for f in idx_h})

    # VMs
    info_rows, ii = rows_from_sheet(wb, "vInfo")
    idx_i = {field: _resolve(ii, aliases) for field, aliases in RVTOOLS_INFO_KEYS.items()}

    vms = []
    for row in info_rows:
        if not row[idx_i["name"]]:
            continue
        if str(row[idx_i["template"]]).lower() == "true":
            continue
        vms.append({
            "name":            row[idx_i["name"]],
            "datacenter":      row[idx_i["datacenter"]],
            "cluster":         row[idx_i["cluster"]],
            "host":            row[idx_i["host"]],
            "cpus":            row[idx_i["cpus"]] or 0,
            "memory_mib":      row[idx_i["memory_mib"]] or 0,
            "provisioned_mib": row[idx_i["provisioned_mib"]] or 0,
            "inuse_mib":       row[idx_i["inuse_mib"]] or 0,
        })

    return hosts, vms, "RVTools", None  # RVTools derives used_mib from per-VM inuse_mib


# ---------------------------------------------------------------------------
# Nutanix Collector loader
# ---------------------------------------------------------------------------
def load_nutanix(wb):
    # --- Build host→datacenter lookup from vCPU sheet ---
    cpu_rows, ci = rows_from_sheet(wb, "vCPU")
    host_to_dc = {}
    for row in cpu_rows:
        hn = get(row, ci, "Host Name")
        dc = get(row, ci, "Datacenter Name")
        if hn and dc:
            host_to_dc[hn] = dc

    # --- Aggregate vCPUs and vRAM per host from vCPU / vMemory sheets ---
    vcpus_by_host  = defaultdict(int)
    for row in cpu_rows:
        hn  = get(row, ci, "Host Name")
        tmpl = str(get(row, ci, "Template") or "").lower()
        if hn and tmpl != "true":
            vcpus_by_host[hn] += get(row, ci, "vCPUs") or 0

    mem_rows, mi = rows_from_sheet(wb, "vMemory")
    vram_by_host = defaultdict(float)
    for row in mem_rows:
        hn   = get(row, mi, "Host Name")
        tmpl = str(get(row, mi, "Template") or "").lower()
        if hn and tmpl != "true":
            vram_by_host[hn] += get(row, mi, "Size (MiB)") or 0

    # --- Hosts ---
    host_rows, hi = rows_from_sheet(wb, "vHosts")
    hosts = []
    for row in host_rows:
        name = get(row, hi, "Host Name")
        if not name:
            continue
        memory_gb = get(row, hi, "Memory Size") or 0   # Nutanix stores GB
        hosts.append({
            "name":             name,
            "datacenter":       host_to_dc.get(name, "Unknown"),
            "cluster":          get(row, hi, "Cluster Name") or "",
            "cpu_model":        get(row, hi, "CPU Model") or "",
            "cpu_mhz":          get(row, hi, "CPU Speed") or 0,
            "num_sockets":      get(row, hi, "CPUs") or 0,
            "cores_per_socket": get(row, hi, "Cores per CPU") or 0,
            "total_cores":      get(row, hi, "CPU Cores") or 0,
            "memory_mb":        float(memory_gb) * 1024,  # convert GB → MB
            "vms_total":        get(row, hi, "VMs") or 0,
            "vms_on":           get(row, hi, "VMs") or 0,
            "vcpus":            vcpus_by_host.get(name, 0),
            "vram_mib":         vram_by_host.get(name, 0),
        })

    # --- VMs from vCPU + vMemory + vDisk ---
    # Build per-VM record keyed by VM name
    vm_map = {}

    for row in cpu_rows:
        name = get(row, ci, "VM Name")
        tmpl = str(get(row, ci, "Template") or "").lower()
        if not name or tmpl == "true":
            continue
        vm_map.setdefault(name, {
            "name": name,
            "datacenter": get(row, ci, "Datacenter Name"),
            "cluster":    get(row, ci, "Cluster Name"),
            "host":       get(row, ci, "Host Name"),
            "cpus": 0, "memory_mib": 0,
            "provisioned_mib": 0, "inuse_mib": None,  # None = not available
        })
        vm_map[name]["cpus"] = get(row, ci, "vCPUs") or 0

    for row in mem_rows:
        name = get(row, mi, "VM Name")
        tmpl = str(get(row, mi, "Template") or "").lower()
        if not name or tmpl == "true":
            continue
        vm_map.setdefault(name, {
            "name": name,
            "datacenter": get(row, mi, "Datacenter Name"),
            "cluster":    get(row, mi, "Cluster Name"),
            "host":       get(row, mi, "Host Name"),
            "cpus": 0, "memory_mib": 0,
            "provisioned_mib": 0, "inuse_mib": None,
        })
        vm_map[name]["memory_mib"] = get(row, mi, "Size (MiB)") or 0

    disk_rows, di = rows_from_sheet(wb, "vDisk")
    for row in disk_rows:
        name = get(row, di, "VM Name")
        if not name or name not in vm_map:
            continue
        vm_map[name]["provisioned_mib"] += get(row, di, "Capacity (MiB)") or 0

    # --- Storage Used: sum Datastore Consumed (MiB) per DC ---
    # Each datastore appears once regardless of how many hosts share it,
    # so summing by DC Name gives the correct total without double-counting.
    ds_rows, dsi = rows_from_sheet(wb, "Datastore")
    dc_used_mib = defaultdict(float)
    for row in ds_rows:
        dc  = get(row, dsi, "DC Name")
        consumed = get(row, dsi, "Consumed (MiB)")
        if dc and consumed is not None:
            dc_used_mib[dc] += float(consumed)

    # Stamp each VM's inuse_mib with the DC-level total so aggregate() can sum it.
    # We do this by setting inuse_mib on a sentinel VM per DC rather than per-VM,
    # but the cleanest approach is to pass dc_used_mib through to aggregate directly.
    vms = list(vm_map.values())
    return hosts, vms, "Nutanix Collector", dict(dc_used_mib)


# ---------------------------------------------------------------------------
# Aggregate per-datacenter
# ---------------------------------------------------------------------------
def aggregate(hosts, vms, dc_used_override=None):
    """
    dc_used_override: optional dict {dc_name: used_mib} — when provided (Nutanix),
    the datastore-level consumed total is used instead of summing per-VM inuse_mib.
    """
    dcs = sorted(set(h["datacenter"] for h in hosts if h["datacenter"]))

    dc_hosts = defaultdict(list)
    for h in hosts:
        dc_hosts[h["datacenter"]].append(h)

    dc_vms = defaultdict(list)
    for v in vms:
        if v.get("datacenter"):
            dc_vms[v["datacenter"]].append(v)

    dc_summary = {}
    for dc in dcs:
        dc_vm_list   = dc_vms[dc]
        dc_host_list = dc_hosts[dc]

        total_vcpus     = sum(v["cpus"] for v in dc_vm_list)
        total_vmem_mib  = sum(v["memory_mib"] for v in dc_vm_list)
        total_alloc_mib = sum(v["provisioned_mib"] for v in dc_vm_list)
        total_pcores    = sum(h["total_cores"] for h in dc_host_list)

        if dc_used_override is not None:
            # Use datastore-level consumed totals (Nutanix path)
            total_used_mib = dc_used_override.get(dc)
        else:
            # Derive from per-VM inuse_mib (RVTools path)
            inuse_values = [v["inuse_mib"] for v in dc_vm_list]
            total_used_mib = (None if any(x is None for x in inuse_values)
                              else sum(inuse_values))

        dc_summary[dc] = {
            "hosts":               dc_host_list,
            "vm_count":            len(dc_vm_list),
            "total_vcpus":         total_vcpus,
            "total_vmem_mib":      total_vmem_mib,
            "alloc_mib":           total_alloc_mib,
            "used_mib":            total_used_mib,
            "total_pcores":        total_pcores,
            "used_from_datastore": dc_used_override is not None,
        }

    return dcs, dc_summary


# ---------------------------------------------------------------------------
# PowerPoint helpers
# ---------------------------------------------------------------------------
def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BG
    bar.line.fill.background()
    add_text_box(slide, title,
                 Inches(0.35), Inches(0.08), Inches(10), Inches(0.65),
                 font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.35), Inches(0.68), Inches(12), Inches(0.38),
                     font_size=13, color=MID_GRAY)


def add_accent_line(slide, top_inches):
    line = slide.shapes.add_shape(1, Inches(0), Inches(top_inches), SLIDE_W, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def cell_set(cell, text, font_size=10, bold=False,
             bg_color=None, font_color=DARK_TEXT, align=PP_ALIGN.LEFT):
    cell.text = str(text) if text is not None else ""
    tf = cell.text_frame
    tf.word_wrap = False
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = font_color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def add_divider_slide(prs, dc_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    stripe = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    add_text_box(slide, dc_name.upper(),
                 Inches(0.5), Inches(2.8), Inches(12), Inches(1.5),
                 font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Datacenter Report",
                 Inches(0.5), Inches(4.1), Inches(12), Inches(0.6),
                 font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)


def add_hosts_slide(prs, dc_name, host_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GRAY)
    add_header_bar(slide, f"{dc_name}  —  ESXi Hosts", f"{len(host_list)} host(s)")
    add_accent_line(slide, 1.1)

    COL_DEFS = [
        ("Host",            2.5,  lambda h: shorten_host(h["name"])),
        ("Cluster",         1.8,  lambda h: h["cluster"] or ""),
        ("CPU Model",       3.2,  lambda h: h["cpu_model"]),
        ("Sockets × Cores", 1.4,  lambda h: f"{h['num_sockets']} × {h['cores_per_socket']}"),
        ("Total Cores",     1.0,  lambda h: h["total_cores"]),
        ("Freq (GHz)",      0.85, lambda h: f"{h['cpu_mhz']/1000:.2f}" if h['cpu_mhz'] else "—"),
        ("Memory (GB)",     1.0,  lambda h: mb_to_gib(h["memory_mb"])),
        ("VMs (total/on)",  1.1,  lambda h: f"{h['vms_total']} / {h['vms_on']}"),
        ("vCPUs",           0.75, lambda h: h["vcpus"] or "—"),
    ]

    n_rows = len(host_list) + 1
    n_cols = len(COL_DEFS)
    table_top    = Inches(1.25)
    table_left   = Inches(0.25)
    table_width  = SLIDE_W - Inches(0.5)
    row_height   = Inches(0.38)
    table_height = row_height * n_rows

    tbl = slide.shapes.add_table(
        n_rows, n_cols, table_left, table_top, table_width, table_height
    ).table

    total_w = sum(w for _, w, _ in COL_DEFS)
    scale = (table_width / Inches(1)) / total_w
    for ci, (_, w, _) in enumerate(COL_DEFS):
        tbl.columns[ci].width = Inches(w * scale)

    for ci, (hdr, _, _) in enumerate(COL_DEFS):
        cell_set(tbl.cell(0, ci), hdr, font_size=10, bold=True,
                 bg_color=DARK_BG, font_color=WHITE, align=PP_ALIGN.CENTER)

    for ri, host in enumerate(host_list, start=1):
        bg = WHITE if ri % 2 == 1 else ROW_ALT
        for ci, (_, _, fn) in enumerate(COL_DEFS):
            align = PP_ALIGN.LEFT if ci <= 2 else PP_ALIGN.CENTER
            cell_set(tbl.cell(ri, ci), fn(host), font_size=9,
                     bg_color=bg, font_color=DARK_TEXT, align=align)

    total_cores = sum(h["total_cores"] for h in host_list)
    total_mem   = sum(mb_to_gib(h["memory_mb"]) for h in host_list)
    total_vms   = sum(h["vms_total"] for h in host_list)
    total_vcpus = sum(h["vcpus"] for h in host_list if h["vcpus"])

    summary = (f"Totals  —  Physical Cores: {total_cores}   "
               f"Memory: {total_mem:.1f} GB   "
               f"VMs (total): {total_vms}   "
               f"vCPUs assigned: {total_vcpus}")
    footer_top = table_top + table_height + Inches(0.12)
    add_text_box(slide, summary, Inches(0.25), footer_top,
                 SLIDE_W - Inches(0.5), Inches(0.35),
                 font_size=10, bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)


def add_summary_slide(prs, dc_name, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GRAY)
    add_header_bar(slide, f"{dc_name}  —  VM Resource Summary",
                   f"{summary['vm_count']} virtual machine(s) (non-template)")
    add_accent_line(slide, 1.1)

    alloc_mib = summary["alloc_mib"]
    used_mib  = summary["used_mib"]   # may be None for Nutanix
    vmem_gib  = mib_to_gib(summary["total_vmem_mib"])
    vcpus     = summary["total_vcpus"]
    pcores    = summary["total_pcores"]

    used_from_datastore = summary.get("used_from_datastore", False)
    if used_mib is not None:
        used_val = f"{mib_to_tib(used_mib)} TiB"
        if used_from_datastore:
            used_sub = "total datastore consumed (incl. vSAN overhead, snapshots, swap)"
        else:
            used_sub = f"{(used_mib/alloc_mib*100 if alloc_mib else 0):.1f}% of allocated"
    else:
        used_val = "N/A"
        used_sub = "not reported by source tool"

    cards = [
        ("Storage Allocated (All VMs)",   f"{mib_to_tib(alloc_mib)} TiB",
         f"({mib_to_gib(alloc_mib):.1f} GiB)"),
        ("Storage Used (All VMs)",         used_val, used_sub),
        ("Memory Allocated (All VMs)",     f"{vmem_gib:.1f} GiB",
         f"({summary['total_vmem_mib']:,} MiB)"),
        ("vCPUs Assigned (All VMs)",       f"{vcpus:,}",
         f"across {summary['vm_count']} VMs"),
        ("Physical CPU Cores (All Hosts)", f"{pcores:,}",
         f"across {len(summary['hosts'])} host(s)"),
        ("pCPU : vCPU Ratio",              ratio_str(pcores, vcpus),
         "physical cores to virtual CPUs"),
    ]

    card_w  = Inches(5.9)
    card_h  = Inches(1.55)
    col_gap = Inches(0.4)
    row_gap = Inches(0.22)
    start_x = [Inches(0.35), Inches(0.35) + card_w + col_gap]
    start_y = Inches(1.3)

    for idx, (label, value, sub) in enumerate(cards):
        col = idx % 2
        row = idx // 2
        x = start_x[col]
        y = start_y + row * (card_h + row_gap)

        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = ACCENT
        card.line.width = Pt(1.5)

        bar = slide.shapes.add_shape(1, x, y, Inches(0.08), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        add_text_box(slide, label,
                     x + Inches(0.18), y + Inches(0.1),
                     card_w - Inches(0.25), Inches(0.38),
                     font_size=11, color=MID_GRAY)
        add_text_box(slide, value,
                     x + Inches(0.18), y + Inches(0.42),
                     card_w - Inches(0.25), Inches(0.72),
                     font_size=28, bold=True, color=DARK_TEXT)
        add_text_box(slide, sub,
                     x + Inches(0.18), y + Inches(1.1),
                     card_w - Inches(0.25), Inches(0.38),
                     font_size=10, color=MID_GRAY, italic=True)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build_presentation(xlsx_path, pptx_path):
    print(f"Reading: {xlsx_path}")
    wb = openpyxl.load_workbook(xlsx_path, read_only=False, data_only=True)

    fmt = detect_format(wb)
    print(f"Format:  {fmt}")

    if fmt == "rvtools":
        hosts, vms, source_label, dc_used_override = load_rvtools(wb)
    else:
        hosts, vms, source_label, dc_used_override = load_nutanix(wb)
    wb.close()

    dcs, dc_summary = aggregate(hosts, vms, dc_used_override)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(cover, DARK_BG)
    stripe = cover.shapes.add_shape(1, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    add_text_box(cover, "Infrastructure Report",
                 Inches(0.5), Inches(1.8), Inches(12), Inches(1.1),
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(cover, f"{len(dcs)} Datacenter(s)  ·  {len(hosts)} ESXi Host(s)  ·  {len(vms)} VM(s)",
                 Inches(0.5), Inches(3.4), Inches(12), Inches(0.5),
                 font_size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(cover, f"Generated from {source_label} export",
                 Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Per-datacenter slides
    for dc in dcs:
        s = dc_summary[dc]
        add_divider_slide(prs, dc)
        add_hosts_slide(prs, dc, s["hosts"])
        add_summary_slide(prs, dc, s)

    prs.save(pptx_path)
    print(f"Saved:   {pptx_path}")
    print(f"Slides:  {len(prs.slides)}  ({len(dcs)} DCs × 3 slides + 1 cover)")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_vcenter_report.py <input.xlsx> [output.pptx]")
        sys.exit(1)
    xlsx = Path(sys.argv[1])
    pptx = Path(sys.argv[2]) if len(sys.argv) > 2 else xlsx.with_suffix(".pptx")
    build_presentation(str(xlsx), str(pptx))
